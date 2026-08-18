"""
Build supervisor meeting slides for Bonolo Masima.
Run from the project folder with the contact_graspnet conda env.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ── Colour palette (UofG-inspired) ──────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x31, 0x5C)   # dark navy
TEAL   = RGBColor(0x00, 0x78, 0x7C)   # UofG teal
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY = RGBColor(0xF2, 0xF2, 0xF2)
DKGRAY = RGBColor(0x44, 0x44, 0x44)
AMBER  = RGBColor(0xFF, 0xA5, 0x00)
GREEN  = RGBColor(0x2E, 0x86, 0x48)
RED    = RGBColor(0xC0, 0x39, 0x2B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helper utilities ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_para(tf, text, size=16, bold=False, color=DKGRAY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=""):
    """Dark navy top bar with title."""
    rect(slide, 0, 0, W, Inches(1.35), fill=NAVY)
    txbox(slide, title,
          Inches(0.35), Inches(0.12), Inches(11), Inches(0.72),
          size=30, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.35), Inches(0.82), Inches(11), Inches(0.45),
              size=16, color=RGBColor(0xAA, 0xCC, 0xDD))


def footer(slide, text="Bonolo Masima · MSc Robotics & AI · University of Glasgow · Jul 2026"):
    rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=TEAL)
    txbox(slide, text,
          Inches(0.3), Inches(7.12), Inches(10), Inches(0.3),
          size=10, color=WHITE)


def bullet_frame(slide, items, x, y, w, h,
                 size=17, color=DKGRAY, indent_char="•  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = indent_char + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_before = Pt(4)
    return tb


def colored_cell(table, row, col, text, bg, fg=WHITE, size=14, bold=False, align=PP_ALIGN.CENTER):
    cell = table.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = fg


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)

    # Teal accent strip
    rect(sl, 0, Inches(5.6), W, Inches(0.12), fill=TEAL)

    txbox(sl,
          "Causal Inference for Robotic Grasp\nFailure Diagnosis under Perceptual Degradation",
          Inches(0.6), Inches(1.3), Inches(12), Inches(2.2),
          size=38, bold=True, color=WHITE)

    txbox(sl, "Supervisor Meeting 3  ·  Thursday 2 July 2026",
          Inches(0.6), Inches(3.7), Inches(10), Inches(0.5),
          size=20, color=TEAL)

    txbox(sl, "Bonolo Masima",
          Inches(0.6), Inches(4.35), Inches(8), Inches(0.5),
          size=22, bold=True, color=WHITE)

    txbox(sl, "MSc Robotics & AI  ·  University of Glasgow",
          Inches(0.6), Inches(4.9), Inches(8), Inches(0.45),
          size=17, color=RGBColor(0xAA, 0xCC, 0xDD))

    txbox(sl, "Research Question:\n"
              "Can a Structural Causal Model diagnose the root cause of a\n"
              "perception-induced grasp failure more accurately than a zero-shot LLM?",
          Inches(0.6), Inches(5.85), Inches(12), Inches(1.4),
          size=15, color=RGBColor(0xCC, 0xEE, 0xFF), italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — 4-Phase project overview
# ════════════════════════════════════════════════════════════════════════════
def slide_phases(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Project Overview: Four-Phase Roadmap")
    footer(sl)

    phases = [
        ("Phase 1", "Simulation\nPipeline", "COMPLETE", TEAL,
         "MuJoCo scene  •  Contact-GraspNet\nIK controller  •  432-trial experiment"),
        ("Phase 2", "Data\nCollection", "COMPLETE", TEAL,
         "Factorial grid (4×4×3×3×3)\n432 trials  •  CSV logged"),
        ("Phase 3", "Structural\nCausal Model", "IN PROGRESS", AMBER,
         "Fit SCM on collected data\nCounterfactual diagnosis"),
        ("Phase 4", "LLM\nBaseline", "UPCOMING", DKGRAY,
         "Zero-shot GPT-4o prompting\nComparison with SCM accuracy"),
    ]

    box_w = Inches(2.8)
    box_h = Inches(4.4)
    gap   = Inches(0.33)
    start_x = Inches(0.5)
    top_y = Inches(1.6)

    for i, (label, title, status, col, desc) in enumerate(phases):
        x = start_x + i * (box_w + gap)

        # Card background
        rect(sl, x, top_y, box_w, box_h, fill=LTGRAY, line=col)

        # Phase label strip
        strip = rect(sl, x, top_y, box_w, Inches(0.42), fill=col)

        txbox(sl, label, x + Inches(0.08), top_y + Inches(0.04),
              box_w - Inches(0.1), Inches(0.35),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txbox(sl, title,
              x + Inches(0.1), top_y + Inches(0.52),
              box_w - Inches(0.2), Inches(0.9),
              size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # Status badge
        badge_col = GREEN if status == "COMPLETE" else (AMBER if status == "IN PROGRESS" else DKGRAY)
        rect(sl, x + Inches(0.3), top_y + Inches(1.55),
             box_w - Inches(0.6), Inches(0.38), fill=badge_col)
        txbox(sl, status,
              x + Inches(0.3), top_y + Inches(1.57),
              box_w - Inches(0.6), Inches(0.34),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txbox(sl, desc,
              x + Inches(0.15), top_y + Inches(2.1),
              box_w - Inches(0.3), Inches(2.0),
              size=14, color=DKGRAY, align=PP_ALIGN.CENTER, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The pipeline (Phase 1)
# ════════════════════════════════════════════════════════════════════════════
def slide_pipeline(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Phase 1: The Simulation Pipeline", "How perception drives grasping")
    footer(sl)

    steps = [
        ("1. Render Depth", "MuJoCo camera\ncaptures RGBD\nfrom controlled\nviewpoint (φ, θ)"),
        ("2. Perturb\nPerception", "Add depth noise σ_d\nDownsample to ρ%\n→ degraded point\ncloud C_pc"),
        ("3. Propose\nGrasps (CGN)", "Contact-GraspNet\nreturns N poses\nwith confidence\nscores q_grasp"),
        ("4. Execute via IK", "Best pose → DLS\nJacobian IK →\narm moves to\ngrasp position"),
        ("5. Log Outcome", "Proximity check:\ne_pose < 6.5 cm?\n→ Y = 1 (success)\nor Y = 0 (fail)"),
    ]

    box_w  = Inches(2.1)
    box_h  = Inches(3.8)
    arr_w  = Inches(0.25)
    gap    = Inches(0.06)
    total  = len(steps) * box_w + (len(steps)-1) * (arr_w + 2*gap)
    start_x = (W - total) / 2
    top_y  = Inches(1.7)

    for i, (title, desc) in enumerate(steps):
        x = start_x + i * (box_w + arr_w + 2*gap)

        rect(sl, x, top_y, box_w, box_h, fill=NAVY)

        txbox(sl, title, x + Inches(0.08), top_y + Inches(0.15),
              box_w - Inches(0.16), Inches(0.85),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(sl, x + Inches(0.35), top_y + Inches(0.95),
             box_w - Inches(0.7), Inches(0.03), fill=TEAL)

        txbox(sl, desc, x + Inches(0.1), top_y + Inches(1.1),
              box_w - Inches(0.2), Inches(2.5),
              size=14, color=RGBColor(0xCC, 0xEE, 0xFF),
              align=PP_ALIGN.CENTER, wrap=True)

        # Arrow between boxes
        if i < len(steps) - 1:
            ax = x + box_w + gap
            ay = top_y + box_h/2 - Inches(0.15)
            txbox(sl, "→", ax, ay, arr_w, Inches(0.4),
                  size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    txbox(sl,
          "Key insight: every component of this pipeline is parameterised by"
          " the exogenous variables — degrading perception cascades to grasp failure.",
          Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.6),
          size=14, color=DKGRAY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Variables
# ════════════════════════════════════════════════════════════════════════════
def slide_variables(prs):
    sl = blank_slide(prs)
    header_bar(sl, "The Causal Variables", "What we measure and why")
    footer(sl)

    rows = [
        # (var, role, what it measures, why it matters, colour)
        ("σ_d",      "Exogenous", "Gaussian noise added to depth buffer", "Controls sensor degradation — root cause 1", NAVY),
        ("ρ",        "Exogenous", "Fraction of point cloud retained",     "Controls observation sparsity — root cause 2", NAVY),
        ("φ, θ",     "Exogenous", "Camera elevation & azimuth",           "Controls geometric visibility of the object",  NAVY),
        ("C_pc",     "Intermediate", "Object pixels visible in image",    "Perception quality — first downstream effect", TEAL),
        ("q_grasp",  "Intermediate", "CGN confidence score (force-closure quality)", "Key proxy: how reliable is the proposed grasp?", TEAL),
        ("e_pose",   "Intermediate", "Euclidean error: proposed vs true position",   "Accuracy of grasp localisation",            TEAL),
        ("n_grasps", "Intermediate", "Number of valid grasp candidates",  "Richness of the CGN output",                   TEAL),
        ("Y",        "Outcome",   "Binary success (e_pose < 6.5 cm)",     "Did the system correctly locate the object?",  AMBER),
    ]

    col_widths = [Inches(1.1), Inches(1.4), Inches(3.8), Inches(4.6)]
    headers = ["Variable", "Role", "What it measures", "Why it matters"]
    col_x = [Inches(0.18), Inches(1.32), Inches(2.75), Inches(6.6)]
    row_h = Inches(0.55)
    top_y = Inches(1.5)

    # Header row
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        rect(sl, cx, top_y, cw, row_h, fill=NAVY)
        txbox(sl, hdr, cx + Inches(0.05), top_y + Inches(0.1),
              cw - Inches(0.1), row_h - Inches(0.1),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (var, role, what, why, col) in enumerate(rows):
        y = top_y + row_h + i * row_h
        bg = LTGRAY if i % 2 == 0 else WHITE
        data = [var, role, what, why]
        for j, (val, cx, cw) in enumerate(zip(data, col_x, col_widths)):
            rect(sl, cx, y, cw, row_h, fill=bg, line=RGBColor(0xDD,0xDD,0xDD))
            txt_col = col if j == 0 else (col if j == 1 else DKGRAY)
            bold = j < 2
            txbox(sl, val, cx + Inches(0.05), y + Inches(0.08),
                  cw - Inches(0.1), row_h - Inches(0.1),
                  size=13, bold=bold, color=txt_col, align=PP_ALIGN.LEFT if j > 1 else PP_ALIGN.CENTER)

    txbox(sl,
          "Confidence score q_grasp is the natural causal mediator:"
          " it collapses perception quality into a single actionable signal for the SCM.",
          Inches(0.18), Inches(6.68), Inches(13), Inches(0.38),
          size=13, italic=True, color=TEAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Results
# ════════════════════════════════════════════════════════════════════════════
def slide_results(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Results: 432-Trial Factorial Experiment", "Clear monotonic causal degradation")
    footer(sl)

    # ── Left: results table ──────────────────────────────────────────────
    table_data = [
        ["σ_d (noise)", "Mean n_grasps", "Mean q_grasp", "Mean e_pose", "Success rate"],
        ["0.000  (clean)",  "34.2",  "0.242",  "5.0 cm",  "71 %"],
        ["0.005  (mild)",   "24.3",  "0.243",  "7.6 cm",  "43 %"],
        ["0.020  (heavy)",  " 3.2",  "0.202",  "11.1 cm", "11 %"],
        ["0.040  (severe)", " 0.1",  "0.163",  "17.5 cm", " 0 %"],
    ]

    col_w  = [Inches(2.1), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)]
    col_x  = [Inches(0.25)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)
    row_h  = Inches(0.56)
    top_y  = Inches(1.55)

    row_bgs = [NAVY, LTGRAY, WHITE, LTGRAY, WHITE]
    row_fgs = [WHITE, DKGRAY, DKGRAY, DKGRAY, DKGRAY]
    success_cols = [WHITE, GREEN, AMBER, RED, RED]

    for i, row in enumerate(table_data):
        y = top_y + i * row_h
        for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
            bg = row_bgs[i]
            fg = row_fgs[i]
            if i == 0:
                rect(sl, cx, y, cw, row_h, fill=NAVY)
                txbox(sl, val, cx+Inches(0.04), y+Inches(0.1),
                      cw-Inches(0.08), row_h-Inches(0.12),
                      size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            else:
                rect(sl, cx, y, cw, row_h, fill=bg, line=RGBColor(0xDD,0xDD,0xDD))
                col = success_cols[i] if j == 4 else fg
                txbox(sl, val, cx+Inches(0.04), y+Inches(0.1),
                      cw-Inches(0.08), row_h-Inches(0.12),
                      size=14, bold=(j==4), color=col, align=PP_ALIGN.CENTER)

    # ── Right: key callouts ──────────────────────────────────────────────
    rx = Inches(8.6)
    callouts = [
        (GREEN,  "71 % success\nin clean conditions",
                 "Contact-GraspNet works\nwell with full perception"),
        (AMBER,  "43 % with\nmild noise (σ=0.005)",
                 "Early degradation already\nhalves performance"),
        (RED,    "0 % at\nσ_d = 0.040",
                 "Severe noise → zero valid\ngrasps → guaranteed failure"),
    ]
    cy = Inches(1.7)
    for col, headline, note in callouts:
        rect(sl, rx, cy, Inches(4.45), Inches(1.5), fill=LTGRAY, line=col)
        rect(sl, rx, cy, Inches(0.15), Inches(1.5), fill=col)
        txbox(sl, headline,
              rx+Inches(0.25), cy+Inches(0.08),
              Inches(4.1), Inches(0.62),
              size=18, bold=True, color=col)
        txbox(sl, note,
              rx+Inches(0.25), cy+Inches(0.72),
              Inches(4.1), Inches(0.68),
              size=13, color=DKGRAY, italic=True)
        cy += Inches(1.65)

    txbox(sl,
          "n_grasps drops 340x from clean to severe noise."
          "  e_pose grows 3.5x.  Both reflect the same causal root.",
          Inches(0.25), Inches(4.5), Inches(8.1), Inches(0.5),
          size=13, italic=True, color=DKGRAY)

    # ── Bar chart proxy (teal bars, proportional) ───────────────────────
    txbox(sl, "Success rate by noise level:", Inches(0.25), Inches(5.1), Inches(5), Inches(0.35),
          size=13, bold=True, color=NAVY)

    bars = [(0.71, GREEN, "σ_d=0"), (0.43, AMBER, "σ_d=0.005"),
            (0.11, RED,   "σ_d=0.02"), (0.0, RED,  "σ_d=0.04")]
    bar_max_w = Inches(5.5)
    bx = Inches(0.25)
    by = Inches(5.55)
    bh = Inches(0.3)
    bgap = Inches(0.12)
    for val, col, lbl in bars:
        w = max(bar_max_w * val, Inches(0.04))
        rect(sl, bx + Inches(1.2), by, w, bh, fill=col)
        txbox(sl, lbl, bx, by, Inches(1.15), bh, size=12, color=DKGRAY)
        txbox(sl, f"{int(val*100)} %",
              bx + Inches(1.2) + w + Inches(0.05), by, Inches(0.6), bh,
              size=12, bold=True, color=col)
        by += bh + bgap


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Defending confidence scores
# ════════════════════════════════════════════════════════════════════════════
def slide_confidence(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Why Use q_grasp (Confidence Score)?", "Defending the key intermediate variable")
    footer(sl)

    # Left panel — what it is
    rect(sl, Inches(0.25), Inches(1.5), Inches(6.1), Inches(5.3), fill=LTGRAY)
    txbox(sl, "What is q_grasp?",
          Inches(0.4), Inches(1.6), Inches(5.8), Inches(0.5),
          size=20, bold=True, color=NAVY)

    points_l = [
        "Learned force-closure quality score from Contact-GraspNet",
        "Trained on 17 million simulated grasps across diverse objects",
        "Outputs a score in [0, 1] — higher = more likely successful grasp",
        "Encodes: finger alignment, approach direction, contact geometry",
        "Not a hand-crafted heuristic — learnt from data",
    ]
    bullet_frame(sl, points_l, Inches(0.4), Inches(2.2), Inches(5.8), Inches(3.5),
                 size=15, color=DKGRAY)

    # Right panel — why use it
    rect(sl, Inches(6.6), Inches(1.5), Inches(6.5), Inches(5.3), fill=NAVY)
    txbox(sl, "Why it fits the SCM",
          Inches(6.75), Inches(1.6), Inches(6.2), Inches(0.5),
          size=20, bold=True, color=WHITE)

    points_r = [
        "Collapses complex perception quality into one scalar — ideal SCM node",
        "q_grasp responds monotonically to perception degradation (σ_d↑ → q↓)",
        "Clean: q_grasp = 0.242 avg,  34 candidates",
        "Noisy (σ_d=0.04): q_grasp = 0.163,  < 1 candidate on average",
        "Satisfies the SCM Markov condition: q_grasp screens off σ_d from Y",
        "Can be used in logistic outcome equation: P(Y=1) = σ(λq_grasp + ...)",
    ]
    bullet_frame(sl, points_r, Inches(6.75), Inches(2.2), Inches(6.1), Inches(4.3),
                 size=15, color=RGBColor(0xCC, 0xEE, 0xFF))

    # Causal equation
    rect(sl, Inches(0.25), Inches(6.62), Inches(12.85), Inches(0.54), fill=TEAL)
    txbox(sl,
          "SCM outcome eq:   Y = f(q_grasp, e_pose, n_grasps, ε)   ←  q_grasp is the natural summary of perception quality",
          Inches(0.4), Inches(6.67), Inches(12.6), Inches(0.42),
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Limitations + improvements
# ════════════════════════════════════════════════════════════════════════════
def slide_limitations(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Current Limitations & Future Improvements")
    footer(sl)

    limit_items = [
        ("Physical lift not achieved",
         "Gripper geometry (primitive capsules) + position-only IK ignores CGN orientation.\n"
         "Fingers approach top-down and miss the cylinder equator."),
        ("Proximity criterion as proxy",
         "Success = e_pose < 6.5 cm (not actual lift). Scientifically defensible for\n"
         "causal analysis, but a real lift would strengthen the thesis."),
        ("Single object type",
         "Only a cylinder tested. Different shapes may change CGN behaviour\n"
         "and should be explored in future work."),
    ]

    future_items = [
        ("6-DoF IK + full CGN orientation",
         "Use the full 4×4 pose matrix from CGN (approach & closing directions)\n"
         "to align the gripper correctly → enables real physical lift"),
        ("Real Panda mesh assets",
         "Switch from primitive capsules to mujoco_menagerie Panda\n"
         "(or robot_descriptions package) for realistic contact geometry"),
        ("SCM counterfactual interventions",
         "Do(σ_d = 0) on a failed trial → did adding noise cause failure?\n"
         "This is the core causal diagnosis comparison with the LLM"),
        ("LLM baseline (Phase 4)",
         "Give GPT-4o the observable variables (C_pc, q_grasp, e_pose, n_grasps, Y)\n"
         "and ask it to identify the root cause — no access to σ_d or ρ"),
    ]

    left_x  = Inches(0.25)
    right_x = Inches(6.8)
    top_y   = Inches(1.6)
    box_h   = Inches(1.42)
    box_w   = Inches(6.3)
    gap     = Inches(0.18)

    # Limitation boxes
    txbox(sl, "Current Limitations", left_x, top_y - Inches(0.38), box_w, Inches(0.35),
          size=16, bold=True, color=RED)
    for i, (title, desc) in enumerate(limit_items):
        y = top_y + i * (box_h + gap)
        rect(sl, left_x, y, box_w, box_h, fill=RGBColor(0xFF,0xEE,0xEE), line=RED)
        rect(sl, left_x, y, Inches(0.15), box_h, fill=RED)
        txbox(sl, title, left_x+Inches(0.24), y+Inches(0.08),
              box_w-Inches(0.3), Inches(0.38),
              size=15, bold=True, color=RED)
        txbox(sl, desc, left_x+Inches(0.24), y+Inches(0.46),
              box_w-Inches(0.3), Inches(0.88),
              size=13, color=DKGRAY, wrap=True)

    # Future boxes
    txbox(sl, "Future Improvements", right_x, top_y - Inches(0.38), box_w, Inches(0.35),
          size=16, bold=True, color=GREEN)
    box_h2 = Inches(1.04)
    for i, (title, desc) in enumerate(future_items):
        y = top_y + i * (box_h2 + gap)
        rect(sl, right_x, y, box_w, box_h2, fill=RGBColor(0xE8,0xF8,0xED), line=GREEN)
        rect(sl, right_x, y, Inches(0.15), box_h2, fill=GREEN)
        txbox(sl, title, right_x+Inches(0.24), y+Inches(0.06),
              box_w-Inches(0.3), Inches(0.35),
              size=15, bold=True, color=GREEN)
        txbox(sl, desc, right_x+Inches(0.24), y+Inches(0.42),
              box_w-Inches(0.3), Inches(0.56),
              size=13, color=DKGRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Next steps
# ════════════════════════════════════════════════════════════════════════════
def slide_next_steps(prs):
    sl = blank_slide(prs)
    header_bar(sl, "Next Steps & Timeline")
    footer(sl)

    weeks = [
        ("Week 1\n(now)",   "Fit SCM",
         "Structural equations via OLS on intermediates;\nlogistic regression on Y",
         TEAL),
        ("Week 2",          "Counterfactual\nDiagnosis",
         "Pearl 3-step procedure on held-out\nfailed trials — SCM root cause",
         TEAL),
        ("Week 3",          "LLM Baseline",
         "Zero-shot GPT-4o on same failures;\ncompare diagnosis accuracy vs SCM",
         NAVY),
        ("Week 4",          "Results &\nDiscussion",
         "Tables, Figs, robustness checks;\nbegin writing Chapter 4 & 5",
         NAVY),
        ("Week 5",          "Full Draft\nto Supervisor",
         "Complete dissertation draft submitted\nfor feedback",
         AMBER),
    ]

    box_w  = Inches(2.35)
    box_h  = Inches(4.0)
    gap    = Inches(0.25)
    top_y  = Inches(1.7)
    start_x = Inches(0.3)

    for i, (week, title, desc, col) in enumerate(weeks):
        x = start_x + i * (box_w + gap)

        # Week number strip
        rect(sl, x, top_y, box_w, Inches(0.5), fill=col)
        txbox(sl, week, x+Inches(0.05), top_y+Inches(0.04),
              box_w-Inches(0.1), Inches(0.44),
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(sl, x, top_y+Inches(0.5), box_w, box_h-Inches(0.5), fill=LTGRAY, line=col)

        txbox(sl, title, x+Inches(0.1), top_y+Inches(0.62),
              box_w-Inches(0.2), Inches(0.88),
              size=17, bold=True, color=col, align=PP_ALIGN.CENTER)

        txbox(sl, desc, x+Inches(0.1), top_y+Inches(1.6),
              box_w-Inches(0.2), Inches(2.2),
              size=13, color=DKGRAY, align=PP_ALIGN.CENTER, wrap=True)

    # Submission note
    rect(sl, Inches(0.3), Inches(5.95), Inches(12.7), Inches(0.55), fill=NAVY)
    txbox(sl,
          "Target dissertation submission: mid-August 2026  |  "
          "Causal analysis complete by end of Week 4",
          Inches(0.4), Inches(6.0), Inches(12.5), Inches(0.44),
          size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_title(prs)
    slide_phases(prs)
    slide_pipeline(prs)
    slide_variables(prs)
    slide_results(prs)
    slide_confidence(prs)
    slide_limitations(prs)
    slide_next_steps(prs)

    out = os.path.join(os.path.dirname(__file__),
                       "Supervisor_Meeting_3_Bonolo_Masima.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
