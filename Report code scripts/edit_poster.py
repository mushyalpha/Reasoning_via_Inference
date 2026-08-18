#!/usr/bin/env python3
"""Edit the existing A1 poster while preserving its seven-panel structure.

The source file is never overwritten. The edited poster is saved as:
    Bonolo_Masima_Poster_edited.pptx
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/bonolomasima/Desktop/Reasoning_via_Inference")
SOURCE = ROOT / "Bonolo_Masima_Poster.pptx"
OUTPUT = ROOT / "Bonolo_Masima_Poster_edited.pptx"
FIG = ROOT / "results" / "figures"

NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x0D, 0x7F, 0x7A)
TEAL_LIGHT = RGBColor(0xE9, 0xF5, 0xF4)
BLUE = RGBColor(0x2F, 0x63, 0x8C)
BLUE_LIGHT = RGBColor(0xE9, 0xF0, 0xF6)
AMBER = RGBColor(0xC9, 0x93, 0x16)
AMBER_LIGHT = RGBColor(0xFB, 0xF4, 0xDE)
RED = RGBColor(0xB5, 0x42, 0x2F)
RED_LIGHT = RGBColor(0xF8, 0xEB, 0xE8)
PURPLE = RGBColor(0x69, 0x52, 0x96)
PURPLE_LIGHT = RGBColor(0xF0, 0xEC, 0xF7)
INK = RGBColor(0x25, 0x30, 0x39)
MUTED = RGBColor(0x50, 0x5C, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xA7, 0xB5, 0xC0)


def delete_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def set_text_style(shape, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    """Apply one style to all existing text in a shape."""
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, x, y, w, h, *, size=12.5, color=INK, gap=2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(gap)
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def rounded_card(slide, x, y, w, h, fill, line=LINE, radius_adjust=0.12):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.adjustments[0] = radius_adjust
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.0)
    return shape


def add_badge(slide, text, x, y, w, fill, color=WHITE):
    card = rounded_card(slide, x, y, w, 0.28, fill, fill, 0.22)
    card.text_frame.clear()
    card.text_frame.margin_left = 0
    card.text_frame.margin_right = 0
    card.text_frame.margin_top = 0
    card.text_frame.margin_bottom = 0
    card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = card.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(9.5)
    run.font.bold = True
    run.font.color.rgb = color
    return card


def add_stat(slide, value, label, x, y, w, color):
    rounded_card(slide, x, y, w, 0.70, WHITE, color, 0.13)
    add_text(
        slide,
        value,
        x + 0.04,
        y + 0.05,
        w - 0.08,
        0.28,
        size=17,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        label,
        x + 0.04,
        y + 0.34,
        w - 0.08,
        0.28,
        size=8.6,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        margin=0,
    )


def add_picture_fit(slide, path, x, y, w, h):
    """Contain an image in the target box without distortion."""
    with Image.open(path) as image:
        iw, ih = image.size
    target_aspect = w / h
    image_aspect = iw / ih
    if image_aspect >= target_aspect:
        draw_w = w
        draw_h = w / image_aspect
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * image_aspect
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    return slide.shapes.add_picture(
        str(path),
        Inches(draw_x),
        Inches(draw_y),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )


def add_arrow(slide, x, y, w, h):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.color.rgb = TEAL
    return arrow


def clear_template_content(slide):
    """Keep the original panel backgrounds, title, header, logo and footer."""
    keep = {
        "Google Shape;44;p7",   # right middle panel background
        "Google Shape;57;p7",   # navy header
        "Google Shape;58;p7",   # centre middle panel background
        "Google Shape;64;p7",   # upper wide panel background
        "Google Shape;81;p7",   # lower take-home panel background
        "Google Shape;86;p7",   # title
        "Google Shape;88;p7",   # navy footer
        "Google Shape;92;p7",   # lower wide panel background
        "Google Shape;97;p7",   # middle-left panel background
        "Google Shape;104;p7",  # top-left panel background
        "Google Shape;110;p7",  # bottom-left panel background
        "Google Shape;122;p7",  # University logo
        "Google Shape;123;p7",  # name and GUID
    }
    for shape in list(slide.shapes):
        if shape.name not in keep:
            delete_shape(shape)

    # Remove old pictures and decorative equation bars embedded in kept groups.
    nested_delete = {
        "Google Shape;67;p7",
        "Google Shape;68;p7",
        "Google Shape;95;p7",
        "Google Shape;100;p7",
    }
    for group in slide.shapes:
        if hasattr(group, "shapes"):
            for child in list(group.shapes):
                if child.name in nested_delete:
                    delete_shape(child)


def add_problem_panel(slide):
    add_text(slide, "Problem & research question", 0.75, 5.16, 4.15, 0.38,
             size=18, bold=True, color=NAVY)
    add_text(
        slide,
        "A failed grasp reveals the outcome—but not whether the root cause was "
        "sensing, viewpoint, object geometry, or execution.",
        0.74, 5.58, 4.17, 0.74, size=12.5, color=INK, line_spacing=1.05,
    )
    rounded_card(slide, 0.73, 6.28, 4.18, 0.68, TEAL_LIGHT, TEAL, 0.12)
    add_text(
        slide,
        "Can a mechanistic SCM identify which controlled factor caused this "
        "specific failure—and recognise when no single fix is sufficient?",
        0.86, 6.38, 3.92, 0.46, size=12.0, bold=True, color=TEAL,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0,
    )
    add_picture_fit(slide, FIG / "fig_pregrasp_gate.png", 0.73, 7.02, 1.62, 1.18)
    add_bullets(
        slide,
        [
            "Contact-GraspNet proposes a top-ranked 6-DoF pose.",
            "The open hand may already intersect the object or table.",
            "A fluent post-hoc explanation cannot verify its own cause.",
        ],
        2.43, 7.04, 2.47, 1.14, size=10.3, color=MUTED, gap=2,
    )


def add_engine_panel(slide):
    add_text(
        slide,
        "Method: a simulator-backed causal inference engine",
        5.82, 5.16, 6.85, 0.36, size=18, bold=True, color=NAVY,
    )
    add_picture_fit(
        slide, FIG / "fig_mechanistic_audit_engine.png",
        5.74, 5.55, 6.55, 2.92,
    )
    rounded_card(slide, 12.32, 5.55, 2.43, 2.86, WHITE, BLUE, 0.10)
    add_text(slide, "What is controlled?", 12.48, 5.70, 2.10, 0.28,
             size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "Depth noise σd",
            "Point retention ρ",
            "Elevation φ and azimuth θ",
        ],
        12.48, 6.02, 2.05, 0.80, size=10.5, color=INK, gap=1,
    )
    add_text(slide, "What makes it causal?", 12.48, 6.85, 2.10, 0.28,
             size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_bullets(
        slide,
        [
            "Graph pre-registered from known dataflow",
            "Same seed replayed; one variable changed",
            "Diagnosis checked by re-simulation",
        ],
        12.48, 7.18, 2.05, 1.02, size=10.2, color=INK, gap=1,
    )


def add_experiment_panel(slide):
    add_text(slide, "Physical experiment & success criterion",
             0.80, 8.67, 4.04, 0.48, size=15.5, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_picture_fit(slide, FIG / "fig_three_objects_fg.png",
                    0.70, 9.23, 4.43, 1.58)
    add_text(
        slide,
        "Three geometries expose different contact and pose-selection failure modes.",
        0.82, 10.80, 4.15, 0.32, size=10.5, color=MUTED,
        align=PP_ALIGN.CENTER, margin=0,
    )
    add_picture_fit(slide, FIG / "fig_depth_degradation_scaled.png",
                    0.70, 11.18, 4.43, 0.92)
    add_text(
        slide,
        "Same simulated scene, progressively corrupted depth.",
        0.82, 12.09, 4.15, 0.28, size=10.2, color=MUTED,
        align=PP_ALIGN.CENTER, margin=0,
    )
    add_picture_fit(slide, FIG / "floating_gripper" / "fig_grasp_sequence.png",
                    0.70, 12.44, 4.43, 0.92)
    add_bullets(
        slide,
        [
            "7,560 trials: 3 objects × 7 σd × 2 ρ × 6 φ × 3 θ × 5 seeds.",
            "Unfiltered top-1 pose; collision gate evaluated with the hand open.",
            "Physical outcome: close, lift 15 cm, and hold through a shake.",
            "Three-stage factorisation separates proposal, geometry, and execution.",
        ],
        0.73, 13.47, 4.35, 1.48, size=10.6, color=INK, gap=2,
    )


def add_gate_panel(slide):
    add_text(slide, "Physical result: the geometric gate",
             5.78, 9.08, 4.18, 0.44, size=16, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_badge(slide, "CONFIRMATORY · PHYSICAL Y", 7.25, 9.50, 2.52, TEAL)
    add_stat(slide, "5.0%", "marginal success", 5.78, 9.84, 1.24, TEAL)
    add_stat(slide, "72.8%", "pre-grasp collision", 7.13, 9.84, 1.32, RED)
    add_stat(slide, "23.1%", "hold | gate-pass", 8.56, 9.84, 1.25, AMBER)
    add_picture_fit(slide, FIG / "fig_trial_flow.png",
                    5.76, 10.68, 4.20, 2.04)
    rounded_card(slide, 5.77, 12.82, 4.18, 1.48, WHITE, LINE, 0.08)
    add_text(
        slide,
        "Depth noise moves the two mechanisms in opposite directions:",
        5.91, 12.93, 3.90, 0.30, size=12.0, bold=True, color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    add_bullets(
        slide,
        [
            "More noise makes poses drift away, so fewer collide at the open-hand gate.",
            "But among poses that clear the gate, physical hold collapses sharply.",
            "The pooled success rate hides this geometric/execution reversal.",
        ],
        5.92, 13.27, 3.84, 0.88, size=10.2, color=MUTED, gap=1,
    )


def add_geometry_panel(slide):
    add_text(slide, "Object geometry dominates execution",
             10.62, 9.08, 4.18, 0.40, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_badge(slide, "CONFIRMATORY · PHYSICAL Y", 11.86, 9.46, 2.64, TEAL)
    add_picture_fit(slide, FIG / "fig_heatmap_success_by_object.png",
                    10.55, 9.79, 4.42, 2.04)
    rounded_card(slide, 10.56, 11.90, 4.40, 0.72, AMBER_LIGHT, AMBER, 0.10)
    add_text(
        slide,
        "Post-gate success: box 1.5%  ·  cylinder 35.4%  ·  mustard 35.0%",
        10.73, 12.04, 4.05, 0.38, size=11.6, bold=True, color=AMBER,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0,
    )
    add_picture_fit(slide, FIG / "fig_rank_rescue.png",
                    10.58, 12.72, 2.68, 1.56)
    rounded_card(slide, 13.38, 12.76, 1.46, 1.44, TEAL_LIGHT, TEAL, 0.12)
    add_text(slide, "89%", 13.45, 12.90, 1.32, 0.42, size=24, bold=True,
             color=TEAL, align=PP_ALIGN.CENTER, margin=0)
    add_text(
        slide,
        "top-20 collision-filtered success at one favourable clean cell",
        13.49, 13.34, 1.24, 0.58, size=9.2, color=INK,
        align=PP_ALIGN.CENTER, margin=0,
    )
    add_text(
        slide,
        "Ceiling check only: post-hoc viewpoint and lighter clamp-and-lift criterion.",
        13.47, 13.94, 1.28, 0.22, size=7.6, color=MUTED,
        align=PP_ALIGN.CENTER, margin=0,
    )


def add_counterfactual_panel(slide):
    add_text(slide, "Counterfactual diagnosis",
             0.83, 15.54, 4.02, 0.40, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    add_picture_fit(slide, FIG / "causal_dag.png",
                    0.72, 16.01, 4.38, 2.48)

    box_w = 1.20
    xs = [0.73, 2.03, 3.33]
    titles = ["1. Abduction", "2. Action", "3. Prediction"]
    bodies = [
        "Recover the realised noise terms U.",
        "Replace one equation with do(X=x′).",
        "Propagate forward to counterfactual Yx′.",
    ]
    fills = [BLUE_LIGHT, PURPLE_LIGHT, TEAL_LIGHT]
    lines = [BLUE, PURPLE, TEAL]
    for i, x in enumerate(xs):
        rounded_card(slide, x, 18.68, box_w, 0.86, fills[i], lines[i], 0.11)
        add_text(slide, titles[i], x + 0.04, 18.76, box_w - 0.08, 0.22,
                 size=10.2, bold=True, color=lines[i], align=PP_ALIGN.CENTER,
                 margin=0)
        add_text(slide, bodies[i], x + 0.08, 19.03, box_w - 0.16, 0.38,
                 size=8.6, color=INK, align=PP_ALIGN.CENTER, margin=0)
        if i < 2:
            add_arrow(slide, x + box_w + 0.03, 18.98, 0.24, 0.22)
    add_text(
        slide,
        "A failure is attributed only when the single reset changes the predicted "
        "outcome; otherwise it remains joint or irreducible within the tested reset.",
        0.77, 19.67, 4.30, 0.51, size=10.2, color=MUTED,
        align=PP_ALIGN.CENTER, margin=0,
    )


def add_pilot_panel(slide):
    add_text(slide, "Pilot diagnosis: SCM vs LLM",
             5.82, 15.08, 5.95, 0.38, size=17, bold=True, color=NAVY)
    add_text(slide, "Scored on the proximity outcome—not physical grasp success",
             8.75, 15.13, 3.82, 0.24, size=9.0, color=MUTED,
             align=PP_ALIGN.RIGHT, margin=0)
    add_badge(slide, "PILOT · PROXIMITY Y", 12.84, 15.08, 1.84, PURPLE)
    add_picture_fit(slide, FIG / "fig_attribution_accuracy.png",
                    5.70, 15.52, 4.15, 2.47)
    add_picture_fit(slide, FIG / "fig_irreducibility_map.png",
                    9.90, 15.52, 4.90, 2.44)
    add_text(
        slide,
        "SCM: 42/95 = 44.2%  |  LLM: 30/95 = 31.6%  |  pre-registered 50% bar missed",
        5.90, 18.03, 4.02, 0.24, size=9.7, bold=True, color=PURPLE,
        align=PP_ALIGN.CENTER, margin=0,
    )
    add_text(
        slide,
        "57.2% of 292 pilot failures have no single-variable fix; irreducibility "
        "concentrates at overhead viewpoints.",
        10.12, 18.02, 4.50, 0.30, size=9.6, bold=True, color=MUTED,
        align=PP_ALIGN.CENTER, margin=0,
    )


def add_takehome_panel(slide):
    add_text(slide, "Take-home message", 5.84, 18.86, 2.65, 0.34,
             size=18, bold=True, color=NAVY)
    add_text(
        slide,
        "Failure is not one mechanism. In the physical grid, top-ranked pose "
        "selection creates a dominant geometric gate, depth noise drives the gate "
        "and execution in opposite directions, and object geometry overwhelms "
        "viewpoint effects.",
        5.83, 19.24, 5.20, 0.91, size=11.6, bold=True, color=INK,
        line_spacing=1.02,
    )
    rounded_card(slide, 11.18, 18.96, 3.56, 1.15, TEAL_LIGHT, TEAL, 0.10)
    add_text(
        slide,
        "The mechanistic audit makes those layers explicit and tests diagnoses by "
        "rewinding the same trial. It improves on the LLM in the pilot, but does "
        "not meet the pre-registered 50% accuracy target.",
        11.35, 19.10, 3.22, 0.82, size=10.5, color=TEAL,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0,
    )


def build():
    prs = Presentation(str(SOURCE))
    slide = prs.slides[0]
    clear_template_content(slide)

    # Keep the original title/header structure; update styling only.
    for shape in slide.shapes:
        if shape.name == "Google Shape;86;p7":
            shape.text = (
                "Counterfactual Diagnosis of Robotic Grasp\n"
                "Failures under Perceptual and Geometric\n"
                "Degradation"
            )
            shape.left = Inches(0.42)
            shape.top = Inches(1.95)
            shape.width = Inches(14.72)
            shape.height = Inches(2.38)
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            set_text_style(shape, size=29, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        elif shape.name == "Google Shape;123;p7":
            shape.text = "MASIMA, BONOLO\n\n3175764M"
            set_text_style(shape, size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text(
        slide,
        "MuJoCo  ·  Contact-GraspNet  ·  7,560 physically scored trials  ·  "
        "simulator-backed counterfactual ground truth",
        0.58, 4.55, 14.40, 0.34, size=15, bold=True, color=TEAL,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0,
    )

    add_problem_panel(slide)
    add_engine_panel(slide)
    add_experiment_panel(slide)
    add_gate_panel(slide)
    add_geometry_panel(slide)
    add_counterfactual_panel(slide)
    add_pilot_panel(slide)
    add_takehome_panel(slide)

    add_text(
        slide,
        "Bonolo Masima  ·  MSc Robotics & AI  ·  Supervisor: Dr Dezong Zhao  ·  "
        "James Watt School of Engineering, University of Glasgow",
        0.55, 21.06, 14.45, 0.36, size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0,
    )

    prs.save(str(OUTPUT))
    print(f"Saved edited copy: {OUTPUT}")
    print(f"Original preserved: {SOURCE}")


if __name__ == "__main__":
    build()
