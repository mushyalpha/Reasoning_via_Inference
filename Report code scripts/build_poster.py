"""
Build the A1 conference poster for Bonolo Masima's MSc project, following the
University of Glasgow School of Engineering A1 poster template layout
(navy header with crest + name/GUID box, title/subhead/paragraph, two columns).

Run: python3 build_poster.py
Output: A1_Poster_Bonolo_Masima.pptx  (A1 portrait, 594mm x 841mm)

Edit the STUDENT_NAME / STUDENT_GUID constants below before printing.
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Student details — EDIT THESE ────────────────────────────────────────────
STUDENT_SURNAME_FIRST = "Masima, Bonolo"
STUDENT_GUID = "2xxxxxxxM"   # <-- replace with your real GUID before printing
SUPERVISOR = "Dezong Zhao"

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "results", "figures")
LOGO = os.path.join(HERE, "assets", "poster", "uofg_logo_crop.png")

# ── Palette (matches template navy + UofG teal) ─────────────────────────────
NAVY    = RGBColor(0x00, 0x1F, 0x3D)   # header navy (matches template crop)
NAVY2   = RGBColor(0x00, 0x31, 0x5C)   # slightly lighter navy for accents
TEAL    = RGBColor(0x00, 0x78, 0x7C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LTGRAY  = RGBColor(0xF3, 0xF4, 0xF5)
MIDGRAY = RGBColor(0xDD, 0xDD, 0xDD)
DKGRAY  = RGBColor(0x33, 0x33, 0x33)
AMBER   = RGBColor(0xFF, 0xA5, 0x00)
GREEN   = RGBColor(0x2E, 0x86, 0x48)
RED     = RGBColor(0xC0, 0x39, 0x2B)

# ── A1 portrait canvas ───────────────────────────────────────────────────────
PAGE_W = Mm(594)
PAGE_H = Mm(841)


def new_prs():
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp


def textbox(slide, text, x, y, w, h, size=18, bold=False, italic=False,
            color=DKGRAY, align=PP_ALIGN.LEFT, font="Arial", anchor=None,
            line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def multi_para(slide, x, y, w, h, paras, anchor=None):
    """paras: list of dicts with text/size/bold/color/space_before/bullet."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    first = True
    for item in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = item.get("align", PP_ALIGN.LEFT)
        if item.get("space_before"):
            p.space_before = Pt(item["space_before"])
        if item.get("line_spacing"):
            p.line_spacing = item["line_spacing"]
        run = p.add_run()
        prefix = "•  " if item.get("bullet") else ""
        run.text = prefix + item["text"]
        run.font.size = Pt(item.get("size", 15))
        run.font.bold = item.get("bold", False)
        run.font.italic = item.get("italic", False)
        run.font.color.rgb = item.get("color", DKGRAY)
        run.font.name = "Arial"
    return tb


def section_header(slide, text, x, y, w, color=NAVY2):
    rect(slide, x, y, Mm(6), Mm(9), fill=color)
    textbox(slide, text, x + Mm(9), y - Mm(1.5), w - Mm(9), Mm(12),
            size=24, bold=True, color=color)
    return y + Mm(13)


def picture_fit(slide, path, x, y, max_w, max_h, caption=None):
    """Insert an image scaled to fit inside (max_w, max_h), centred, with an
    optional caption below. Returns the y-coordinate just below the block."""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)
    px = x + (max_w - w) // 2
    slide.shapes.add_picture(path, px, y, width=w, height=h)
    bottom = y + h
    if caption:
        textbox(slide, caption, x, bottom + Mm(1), max_w, Mm(10),
                size=12.5, italic=True, color=DKGRAY, align=PP_ALIGN.CENTER)
        bottom += Mm(10)
    return bottom + Mm(4)


# ════════════════════════════════════════════════════════════════════════════
def build():
    prs = new_prs()
    sl = blank_slide(prs)

    margin = Mm(18)
    content_w = PAGE_W - 2 * margin

    # ── Header bar ───────────────────────────────────────────────────────
    header_h = Mm(50)
    rect(sl, 0, 0, PAGE_W, header_h, fill=NAVY)

    sl.shapes.add_picture(LOGO, margin, Mm(15), height=Mm(20))

    name_box_w = Mm(150)
    name_box_x = PAGE_W - margin - name_box_w
    rect(sl, name_box_x, Mm(10), name_box_w, Mm(30), fill=NAVY2)
    textbox(sl, STUDENT_SURNAME_FIRST, name_box_x + Mm(6), Mm(13),
            name_box_w - Mm(12), Mm(12), size=20, bold=True, color=WHITE)
    textbox(sl, f"GUID: {STUDENT_GUID}", name_box_x + Mm(6), Mm(27),
            name_box_w - Mm(12), Mm(10), size=16, bold=True, color=WHITE)

    # ── Title / subhead / intro paragraph ───────────────────────────────
    y = header_h + Mm(10)
    textbox(sl, "Causal Inference for Robotic Grasp Failure Diagnosis "
                "under Perceptual Degradation",
            margin, y, content_w, Mm(45), size=58, bold=True, color=NAVY2,
            line_spacing=1.02)
    y += Mm(46)

    textbox(sl, "Can a Structural Causal Model diagnose why a robot grasp "
                "failed more reliably than a zero-shot LLM?",
            margin, y, content_w, Mm(16), size=32, bold=True, color=TEAL)
    y += Mm(17)

    intro = ("Robots frequently fail to grasp objects when perception is "
             "imperfect \u2014 noisy depth, sparse point clouds, or poor "
             "viewpoints \u2014 but knowing that a grasp failed is not the "
             "same as knowing why. This project builds a Structural Causal "
             "Model (SCM) that identifies which controlled perceptual "
             "perturbation (depth noise \u03c3\u1d05, sparsity \u03c1, camera "
             "elevation \u03c6, azimuth \u03b8) caused a Contact-GraspNet "
             "grasp to fail, using Pearl's counterfactual (abduction\u2192"
             "action\u2192prediction) reasoning \u2014 and benchmarks its "
             "diagnostic accuracy against a zero-shot LLM given the same "
             "evidence.")
    textbox(sl, intro, margin, y, content_w, Mm(30), size=24, color=DKGRAY,
            line_spacing=1.08)
    y += Mm(34)

    col_gap = Mm(12)
    col_w = (content_w - col_gap) / 2
    col1_x = margin
    col2_x = margin + col_w + col_gap
    col_top = y

    # ══════════════════════════════════════════════════════════════════
    # COLUMN 1 — Method & Causal Pipeline
    # ══════════════════════════════════════════════════════════════════
    cy = section_header(sl, "Column 1 \u2014 Method & Causal Pipeline", col1_x, col_top, col_w)
    cy += Mm(4)

    # Pipeline flow (drawn as boxes + arrows, no image asset needed)
    steps = [
        ("1. Render", "MuJoCo camera\ncaptures depth from\nviewpoint (\u03c6,\u03b8)"),
        ("2. Degrade", "Inject noise \u03c3\u1d05\nDownsample \u03c1\n\u2192 point cloud C_pc"),
        ("3. Propose", "Contact-GraspNet\nreturns grasps +\nconfidence q_grasp"),
        ("4. Execute", "Floating-gripper\nshake test on\ntop-1 pose"),
        ("5. Outcome", "Success / failure Y\nlogged to CSV"),
    ]
    n = len(steps)
    arrow_w = Mm(6)
    box_w = (col_w - (n - 1) * arrow_w) / n
    box_h = Mm(38)
    bx = col1_x
    for i, (title, desc) in enumerate(steps):
        rect(sl, bx, cy, box_w, box_h, fill=NAVY2)
        textbox(sl, title, bx + Mm(2), cy + Mm(2), box_w - Mm(4), Mm(10),
                size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx + Mm(3), cy + Mm(12), box_w - Mm(6), Mm(0.6), fill=TEAL)
        textbox(sl, desc, bx + Mm(1.5), cy + Mm(13.5), box_w - Mm(3), Mm(24),
                size=11, color=RGBColor(0xE0, 0xEE, 0xF2), align=PP_ALIGN.CENTER,
                line_spacing=1.0)
        bx += box_w
        if i < n - 1:
            textbox(sl, "\u2192", bx, cy + box_h / 2 - Mm(6), arrow_w, Mm(12),
                    size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
            bx += arrow_w
    cy += box_h + Mm(8)

    # Causal DAG figure
    textbox(sl, "The causal graph (dataflow-derived, dated pre-registration)",
            col1_x, cy, col_w, Mm(8), size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    cy = picture_fit(sl, os.path.join(FIG, "scm_dag_corrected.png"),
                      col1_x, cy, col_w, Mm(95),
                      caption="Fig. 1 \u2014 Corrected causal DAG. Every edge/non-edge tested "
                              "interventionally against the 432-trial dataset "
                              "(test_dag_edges.py) \u2014 all tests passed.")

    # Variable table
    textbox(sl, "Causal variables", col1_x, cy, col_w, Mm(8),
            size=16, bold=True, color=NAVY2)
    cy += Mm(9)

    var_rows = [
        ("\u03c3\u1d05", "Exogenous", "Gaussian depth-buffer noise", "0, 0.005, 0.02, 0.04 m"),
        ("\u03c1", "Exogenous", "Point cloud downsample fraction", "1.0, 0.75, 0.5, 0.25"),
        ("\u03c6, \u03b8", "Exogenous", "Camera elevation / azimuth", "30\u201360\u00b0, 0\u201390\u00b0"),
        ("C_pc", "Mediator", "Point cloud completeness", "measured"),
        ("q_grasp", "Mediator", "CGN top-1 confidence", "measured"),
        ("e_pose", "Mediator", "Grasp pose error vs. oracle", "measured"),
        ("Y", "Outcome", "Grasp success (shake test)", "binary"),
    ]
    headers = ["Var", "Role", "Meaning", "Domain"]
    widths = [Mm(16), Mm(24), col_w - Mm(16) - Mm(24) - Mm(38), Mm(38)]
    row_h = Mm(9)
    rx0 = col1_x
    xs = [rx0]
    for wi in widths[:-1]:
        xs.append(xs[-1] + wi)
    for j, (hd, xx, ww) in enumerate(zip(headers, xs, widths)):
        rect(sl, xx, cy, ww, row_h, fill=NAVY2)
        textbox(sl, hd, xx + Mm(1.5), cy + Mm(1.2), ww - Mm(3), row_h - Mm(2),
                size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cy += row_h
    for i, row in enumerate(var_rows):
        bg = LTGRAY if i % 2 == 0 else WHITE
        for j, (val, xx, ww) in enumerate(zip(row, xs, widths)):
            rect(sl, xx, cy, ww, row_h, fill=bg, line=MIDGRAY, line_w=Pt(0.5))
            textbox(sl, val, xx + Mm(1.5), cy + Mm(1.2), ww - Mm(3), row_h - Mm(2),
                    size=12, bold=(j == 0), color=NAVY2 if j == 0 else DKGRAY,
                    align=PP_ALIGN.CENTER if j < 2 else PP_ALIGN.LEFT)
        cy += row_h
    cy += Mm(6)

    # Counterfactual method
    textbox(sl, "Counterfactual diagnosis procedure (Pearl, 3-step)",
            col1_x, cy, col_w, Mm(8), size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    cf_items = [
        {"text": "Abduction \u2014 infer the exogenous noise term \u03b5 for a failed "
                  "trial from its observed evidence.", "bullet": True, "size": 14.5,
         "space_before": 3},
        {"text": "Action \u2014 intervene do(\u03c3\u1d05 = 0), do(\u03c1 = 1), etc., one "
                  "variable at a time, holding \u03b5 fixed.", "bullet": True, "size": 14.5,
         "space_before": 3},
        {"text": "Prediction \u2014 re-simulate; if success flips to 1, that variable "
                  "is the diagnosed cause.", "bullet": True, "size": 14.5, "space_before": 3},
        {"text": "Applied to all 292 real failed trials \u00d7 4 single-variable "
                  "interventions = 1,168 re-simulations (ground truth).",
         "bullet": True, "size": 14.5, "space_before": 3},
    ]
    multi_para(sl, col1_x, cy, col_w, Mm(40), cf_items)
    cy += Mm(42)

    # SCM structural-equation fit summary
    textbox(sl, "SCM structural equations \u2014 fit quality", col1_x, cy, col_w, Mm(8),
            size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    eq_rows = [
        ("Eq1", "C_pc ~ \u03c6 + \u03b8", "OLS", "R\u00b2 = 0.893"),
        ("Eq2A", "has_grasps ~ \u03c3\u1d05+\u03c1+\u03c6+\u03b8", "Logistic", "pseudo-R\u00b2=0.554, AUC=0.943"),
        ("Eq2B", "n_grasps ~ \u03c3\u1d05+\u03c1+\u03c6+\u03b8", "Neg. Binomial", "\u03c3\u1d05 & \u03c1 both sig."),
        ("Eq3", "q_grasp ~ log(n_grasps)+\u2026", "OLS", "R\u00b2 = 0.699"),
        ("Eq4", "e_pose ~ \u03c3\u1d05+\u03c1+\u2026", "OLS", "\u03c1 90.9% mediated via q_grasp"),
    ]
    eqw = [Mm(16), Mm(90), Mm(38), col_w - Mm(16) - Mm(90) - Mm(38)]
    exs = [col1_x]
    for wi in eqw[:-1]:
        exs.append(exs[-1] + wi)
    ehdrs = ["Eq", "Structural form", "Method", "Fit statistic"]
    row_h2 = Mm(9)
    for hd, xx, ww in zip(ehdrs, exs, eqw):
        rect(sl, xx, cy, ww, row_h2, fill=NAVY2)
        textbox(sl, hd, xx + Mm(1.5), cy + Mm(1.2), ww - Mm(3), row_h2 - Mm(2),
                size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cy += row_h2
    for i, row in enumerate(eq_rows):
        bg = LTGRAY if i % 2 == 0 else WHITE
        for j, (val, xx, ww) in enumerate(zip(row, exs, eqw)):
            rect(sl, xx, cy, ww, row_h2, fill=bg, line=MIDGRAY, line_w=Pt(0.5))
            textbox(sl, val, xx + Mm(1.5), cy + Mm(1.2), ww - Mm(3), row_h2 - Mm(2),
                    size=11.5, bold=(j == 0), color=NAVY2 if j == 0 else DKGRAY,
                    align=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)
        cy += row_h2
    cy += Mm(8)

    cy = picture_fit(sl, os.path.join(FIG, "scm_coefficients.png"),
                      col1_x, cy, col_w, Mm(90),
                      caption="Fig. 1b \u2014 Fitted SCM coefficients per structural equation. "
                              "\u03c3\u1d05 dominates pipeline collapse; \u03c1's effect is almost "
                              "entirely mediated through grasp count / confidence, not a "
                              "direct effect on C_pc.")

    col1_bottom = cy

    # ══════════════════════════════════════════════════════════════════
    # COLUMN 2 — Results & Contribution
    # ══════════════════════════════════════════════════════════════════
    cy = section_header(sl, "Column 2 \u2014 Results & Contribution", col2_x, col_top, col_w)
    cy += Mm(4)

    # Headline stat cards
    cards = [
        (GREEN, "134 / 432", "trials succeeded (31.0%) across the full "
         "factorial grid \u2014 clear monotonic degradation as \u03c3\u1d05\u2191 / "
         "\u03c1\u2193 / \u03c6\u2191."),
        (RED, "57.2%", "of failures (167/292) have no single-variable fix \u2014 "
         "irreducible or multi-causal, mostly at \u03c6=60\u00b0 (overhead "
         "geometry)."),
        (AMBER, "10.3% vs 57.2%", "zero-shot LLM full-set diagnostic accuracy "
         "vs. the trivial \u201calways guess \u2018none\u2019\u201d majority-class "
         "baseline \u2014 the LLM underperforms a naive baseline."),
    ]
    card_h = Mm(28)
    for col, headline, note in cards:
        rect(sl, col2_x, cy, col_w, card_h, fill=LTGRAY, line=col, line_w=Pt(1.5))
        rect(sl, col2_x, cy, Mm(2.5), card_h, fill=col)
        textbox(sl, headline, col2_x + Mm(6), cy + Mm(2), col_w - Mm(10), Mm(10),
                size=22, bold=True, color=col)
        textbox(sl, note, col2_x + Mm(6), cy + Mm(12), col_w - Mm(10), Mm(15),
                size=13, color=DKGRAY, line_spacing=1.0)
        cy += card_h + Mm(5)
    cy += Mm(2)

    # Primary cause breakdown chart
    textbox(sl, "Root cause of failure: single-cause vs. irreducible",
            col2_x, cy, col_w, Mm(8), size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    cy = picture_fit(sl, os.path.join(FIG, "counterfactual_primary_cause_breakdown.png"),
                      col2_x, cy, col_w, Mm(80),
                      caption="Fig. 2 \u2014 Ground-truth primary cause per failed trial, "
                              "from 1,168 counterfactual re-simulations.")

    # SCM fit quality
    textbox(sl, "SCM structural fit", col2_x, cy, col_w, Mm(8),
            size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    cy = picture_fit(sl, os.path.join(FIG, "scm_heatmap_sigma_rho.png"),
                      col2_x, cy, col_w, Mm(75),
                      caption="Fig. 3 \u2014 Success rate over the (\u03c3\u1d05, \u03c1) grid. "
                              "\u03c3\u1d05 dominates (OR\u22480 at high noise); \u03c1 effect is "
                              "largely mediated through grasp count, not confidence.")

    # LLM baseline comparison
    textbox(sl, "LLM baseline: diagnostic accuracy", col2_x, cy, col_w, Mm(8),
            size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    cy = picture_fit(sl, os.path.join(FIG, "llm_baseline_primary_accuracy.png"),
                      col2_x, cy, col_w, Mm(70),
                      caption="Fig. 4 \u2014 Zero-shot LLM per-cause accuracy on the 95 trials "
                              "with a clear single primary cause (overall 31.6%). Note: "
                              "0% on \u03b8 \u2014 azimuth failures are systematically missed.")

    # Limitations / future work box
    textbox(sl, "Limitations & next steps", col2_x, cy, col_w, Mm(8),
            size=16, bold=True, color=NAVY2)
    cy += Mm(9)
    lim_items = [
        {"text": "SCM-vs-LLM head-to-head diagnostic accuracy: SCM counterfactual "
                  "ranking on the 292 failed trials vs. ground truth \u2014 in progress.",
         "bullet": True, "size": 14.5, "space_before": 2},
        {"text": "Single-object (cylinder) scene \u2192 multi-object redesign "
                  "(box + mustard bottle, clutter, occlusion) underway per marker "
                  "feedback.", "bullet": True, "size": 14.5, "space_before": 2},
        {"text": "Recovery actions are out of scope \u2014 diagnosis only.",
         "bullet": True, "size": 14.5, "space_before": 2},
    ]
    multi_para(sl, col2_x, cy, col_w, Mm(35), lim_items)
    cy += Mm(36)

    col2_bottom = cy

    # ── Take-home message banner (fills remaining space before footer) ────
    banner_y = max(col1_bottom, col2_bottom) + Mm(10)
    footer_h = Mm(16)
    footer_y = PAGE_H - Mm(20)
    banner_h = max(footer_y - Mm(10) - banner_y, Mm(45))
    banner_h = min(banner_h, Mm(90))

    rect(sl, margin, banner_y, content_w, banner_h, fill=NAVY)
    rect(sl, margin, banner_y, Mm(3), banner_h, fill=TEAL)
    textbox(sl, "Take-home message", margin + Mm(10), banner_y + Mm(6),
            content_w - Mm(20), Mm(11), size=22, bold=True, color=TEAL)
    textbox(sl,
            "A causal model can say not just that a grasp failed, but why "
            "\u2014 and, honestly, when it can't (57.2% of failures are "
            "irreducible to a single variable). Zero-shot LLMs, given the "
            "same evidence, currently underperform even a naive "
            "majority-class guess. The next step is a direct SCM-vs-LLM "
            "diagnostic accuracy comparison on the same 292 failed trials.",
            margin + Mm(10), banner_y + Mm(19), content_w - Mm(20),
            banner_h - Mm(23), size=17, color=RGBColor(0xE0, 0xEE, 0xF2),
            line_spacing=1.15)

    # ── Footer ───────────────────────────────────────────────────────────
    rect(sl, 0, footer_y, PAGE_W, Mm(0.6), fill=MIDGRAY)
    textbox(sl, f"Bonolo Masima  \u00b7  MSc Robotics & AI  \u00b7  Supervisor: {SUPERVISOR}"
                "  \u00b7  School of Engineering, University of Glasgow",
            margin, footer_y + Mm(3), content_w - Mm(90), Mm(8),
            size=13, color=DKGRAY)
    textbox(sl, "University of Glasgow, charity number SC004401",
            margin, footer_y + Mm(3), content_w, Mm(8),
            size=11, color=DKGRAY, align=PP_ALIGN.RIGHT)

    out = os.path.join(HERE, "A1_Poster_Bonolo_Masima.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Canvas: A1 portrait, 594 x 841 mm")
    print(f"Column bottom (col1): {col1_bottom}")
    print(f"Column bottom (col2): {col2_bottom}")


if __name__ == "__main__":
    build()
